import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
prs = Presentation("Hotel_CI_Business_Case.pptx")
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t[:70])
    print(f"Slide {i}: " + " | ".join(texts[:5]))
print(f"\nTotal slides: {len(prs.slides)}")
